# SPDX-License-Identifier: MIT
# Copyright (c) 2025 Elena Viter

# tools/file_text_extractor.py
from __future__ import annotations

import io
import mimetypes
import zipfile
import importlib
from dataclasses import dataclass, field
from typing import Optional, Tuple, Dict, Any, List

# Optional deps (loaded lazily):
# - PyMuPDF (fitz)
# - pdfminer.six
# - pypdf
# - python-docx
# - python-pptx
# - chardet


@dataclass
class ExtractInfo:
    mime: str
    ext: str
    meta: Dict[str, Any] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)


class DocumentTextExtractor:
    """Extract text from PDF, DOCX, PPTX, TXT/MD. Provide bytes, filename, optional mime."""

    # Minimal ext->MIME map for when mimetypes.guess_type is vague/missing
    _EXT_MIME = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".txt": "text/plain",
        ".md": "text/markdown",
    }

    def extract(self, data: bytes, filename: str, mime: Optional[str] = None) -> Tuple[str, ExtractInfo]:
        mime, ext, hints = self._resolve_mime_and_ext(data, filename, mime)
        text, meta, warnings = "", {}, []

        try:
            if mime == "application/pdf" or ext == ".pdf" or hints.get("is_pdf", False):
                text, meta, more = self._extract_pdf(data)
                warnings.extend(more)

            elif mime.endswith("wordprocessingml.document") or ext == ".docx" or hints.get("is_docx", False):
                text, meta, more = self._extract_docx(data)
                warnings.extend(more)

            elif mime.endswith("presentationml.presentation") or ext == ".pptx" or hints.get("is_pptx", False):
                text, meta, more = self._extract_pptx(data)
                warnings.extend(more)

            elif mime.endswith("spreadsheetml.sheet") or ext == ".xlsx" or hints.get("is_xlsx", False):
                text, meta, more = self._extract_xlsx(data)
                warnings.extend(more)

            elif mime.startswith("text/") or ext in (".txt", ".md"):
                text, meta, more = self._extract_textlike(data, mime or "text/plain")
                warnings.extend(more)

            else:
                # Last-ditch: try decoding as text
                text, meta, more = self._extract_textlike(data, "text/plain")
                warnings.append(f"Unknown MIME '{mime or ''}'. Decoded as text/plain.")
                warnings.extend(more)

        except Exception as e:
            warnings.append(f"Extraction error: {e!r}")

        info = ExtractInfo(mime=mime, ext=ext, meta=meta, warnings=warnings)
        return self._normalize(text), info

    # ---------- MIME / type detection ----------

    def _resolve_mime_and_ext(self, data: bytes, filename: str, mime: Optional[str]) -> Tuple[str, str, Dict[str, Any]]:
        ext = ""
        if "." in filename:
            ext = "." + filename.lower().rsplit(".", 1)[-1]

        guessed = (mimetypes.guess_type(filename)[0] or "").lower()
        hard_map = self._EXT_MIME.get(ext, "")

        # Quick file-signature hints
        hints: Dict[str, Any] = {
            "is_pdf": data.startswith(b"%PDF"),
            "is_zip": data[:2] == b"PK",
        }

        if hints["is_zip"]:
            # Peek inside OOXML to distinguish DOCX vs PPTX
            try:
                with zipfile.ZipFile(io.BytesIO(data)) as z:
                    names = set(z.namelist())
                    if "[Content_Types].xml" in names:
                        ct = z.read("[Content_Types].xml").decode("utf-8", "replace")
                        if "wordprocessingml.document" in ct:
                            hints["is_docx"] = True
                        if "presentationml.presentation" in ct:
                            hints["is_pptx"] = True
                        if "spreadsheetml.sheet" in ct:
                            hints["is_xlsx"] = True
            except Exception:
                pass

        # Decide MIME: explicit > specific guess > hard map > header-based > fallback
        chosen = (mime or "").lower().strip()
        if not chosen or chosen == "application/octet-stream":
            chosen = guessed or hard_map
            if not chosen:
                if hints["is_pdf"]:
                    chosen = "application/pdf"
                elif hints.get("is_docx"):
                    chosen = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                elif hints.get("is_pptx"):
                    chosen = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                elif hints.get("is_xlsx"):
                    chosen = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                elif ext in (".md",):
                    chosen = "text/markdown"
                elif ext in (".txt",):
                    chosen = "text/plain"
                else:
                    chosen = "application/octet-stream"

        # Normalize ext for consistency
        if not ext:
            # derive from MIME if possible
            for k, v in self._EXT_MIME.items():
                if v == chosen:
                    ext = k
                    break
        return chosen, ext, hints

    # ---------- Extractors ----------

    def _extract_pdf(self, data: bytes) -> Tuple[str, Dict[str, Any], List[str]]:
        warnings: List[str] = []
        # Try PyMuPDF (best quality/robustness)
        fitz = importlib.util.find_spec("fitz")
        if fitz:
            try:
                import fitz  # type: ignore
                doc = fitz.open(stream=data, filetype="pdf")
                parts = []
                for i, page in enumerate(doc, start=1):
                    parts.append(f"\n\n--- Page {i} ---\n")
                    # "text" is a decent default; consider "blocks" for layout-sensitive content
                    parts.append(page.get_text("text"))
                return "".join(parts), {"pages": doc.page_count, "engine": "pymupdf"}, warnings
            except Exception as e:
                warnings.append(f"PyMuPDF failed: {e!r}")

        # Fallback: pdfminer.six
        pdfminer = importlib.util.find_spec("pdfminer")
        if pdfminer:
            try:
                from pdfminer.high_level import extract_text  # type: ignore
                text = extract_text(io.BytesIO(data)) or ""
                return text, {"engine": "pdfminer.six"}, warnings
            except Exception as e:
                warnings.append(f"pdfminer failed: {e!r}")

        # Fallback: pypdf
        pypdf = importlib.util.find_spec("pypdf") or importlib.util.find_spec("PyPDF2")
        if pypdf:
            try:
                try:
                    from pypdf import PdfReader  # type: ignore
                except Exception:
                    from PyPDF2 import PdfReader  # type: ignore
                reader = PdfReader(io.BytesIO(data))
                parts = []
                for i, page in enumerate(reader.pages, start=1):
                    parts.append(f"\n\n--- Page {i} ---\n")
                    parts.append(page.extract_text() or "")
                return "".join(parts), {"pages": len(reader.pages), "engine": "pypdf"}, warnings
            except Exception as e:
                warnings.append(f"pypdf failed: {e!r}")

        # Nothing worked
        warnings.append("No PDF extractor available; returning empty text.")
        return "", {"engine": None}, warnings

    def _extract_docx(self, data: bytes) -> Tuple[str, Dict[str, Any], List[str]]:
        warnings: List[str] = []
        spec = importlib.util.find_spec("docx")
        if not spec:
            return "", {"engine": None}, ["python-docx not installed"]

        from docx import Document  # type: ignore

        doc = Document(io.BytesIO(data))
        parts: List[str] = []

        # Paragraphs
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)

        # Tables
        try:
            for tbl in doc.tables:
                parts.append("")  # break
                for row in tbl.rows:
                    cells = [" ".join(c.text.strip().split()) for c in row.cells]
                    parts.append(" | ".join(cells))
        except Exception as e:
            warnings.append(f"Table parse failed: {e!r}")

        return "\n".join(parts), {"engine": "python-docx"}, warnings

    def _extract_pptx(self, data: bytes) -> Tuple[str, Dict[str, Any], List[str]]:
        warnings: List[str] = []
        spec = importlib.util.find_spec("pptx")
        if not spec:
            return "", {"engine": None}, ["python-pptx not installed"]

        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(data))
        parts: List[str] = []
        slide_count = 0

        def shape_text(sh) -> List[str]:
            out: List[str] = []
            try:
                if hasattr(sh, "has_text_frame") and sh.has_text_frame:
                    for p in sh.text_frame.paragraphs:
                        t = " ".join(run.text or "" for run in p.runs).strip()
                        if t:
                            out.append(t)
                # recurse into grouped shapes
                if hasattr(sh, "shapes"):
                    for sub in sh.shapes:
                        out.extend(shape_text(sub))
            except Exception:
                pass
            return out

        for i, slide in enumerate(prs.slides, start=1):
            slide_count += 1
            parts.append(f"\n\n--- Slide {i} ---\n")
            for sh in slide.shapes:
                parts.extend(shape_text(sh))
            # notes (if any)
            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    nt = slide.notes_slide.notes_text_frame.text
                    if nt and nt.strip():
                        parts.append("\n[Notes]\n" + nt.strip())
            except Exception:
                pass

        return "\n".join(parts), {"engine": "python-pptx", "slides": slide_count}, warnings

    def _extract_xlsx(self, data: bytes) -> Tuple[str, Dict[str, Any], List[str]]:
        warnings: List[str] = []
        spec = importlib.util.find_spec("openpyxl")
        if not spec:
            return "", {"engine": None}, ["openpyxl not installed"]

        from openpyxl import load_workbook  # type: ignore

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: List[str] = []
        sheet_count = 0
        for ws in wb.worksheets:
            sheet_count += 1
            parts.append(f"\n\n--- Sheet {ws.title} ---\n")
            try:
                for row in ws.iter_rows(values_only=True):
                    cells = []
                    for v in row:
                        if v is None:
                            cells.append("")
                        else:
                            cells.append(str(v).strip())
                    line = " | ".join(cells).strip()
                    if line:
                        parts.append(line)
            except Exception as e:
                warnings.append(f"Worksheet parse failed ({ws.title}): {e!r}")
        return "\n".join(parts), {"engine": "openpyxl", "sheets": sheet_count}, warnings

    def _extract_textlike(self, data: bytes, mime: str) -> Tuple[str, Dict[str, Any], List[str]]:
        # Try utf-8; fall back to chardet
        warnings: List[str] = []
        try:
            text = data.decode("utf-8")
            enc = "utf-8"
        except UnicodeDecodeError:
            chardet_spec = importlib.util.find_spec("chardet")
            if chardet_spec:
                import chardet  # type: ignore
                det = chardet.detect(data) or {}
                enc = (det.get("encoding") or "utf-8")
                try:
                    text = data.decode(enc, errors="replace")
                    if det.get("confidence", 0) < 0.7:
                        warnings.append(f"Low-confidence encoding detection ({det.get('confidence')}). Used {enc}.")
                except Exception as e:
                    warnings.append(f"Decoding with {enc} failed: {e!r}")
                    text = data.decode("utf-8", errors="replace")
                    enc = "utf-8"
            else:
                warnings.append("chardet not installed; forced utf-8 with replacement.")
                text = data.decode("utf-8", errors="replace")
                enc = "utf-8"

        return text, {"engine": "decode", "encoding": enc, "mime": mime}, warnings

    # ---------- utilities ----------

    @staticmethod
    def _normalize(s: str) -> str:
        # Normalize newlines; trim overly long runs of blank lines
        s = s.replace("\r\n", "\n").replace("\r", "\n")
        # collapse >2 blank lines to just 2
        out_lines: List[str] = []
        blank_run = 0
        for line in s.split("\n"):
            if line.strip() == "":
                blank_run += 1
                if blank_run <= 2:
                    out_lines.append("")
            else:
                blank_run = 0
                out_lines.append(line.rstrip())
        return "\n".join(out_lines)
