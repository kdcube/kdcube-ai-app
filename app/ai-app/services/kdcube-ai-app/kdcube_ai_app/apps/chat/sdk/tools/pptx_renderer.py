# SPDX-License-Identifier: MIT
# Copyright (c) 2025 Elena Viter
# chat/sdk/tools/pptx_renderer.py

from __future__ import annotations
import pathlib
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
import json
import re
from html.parser import HTMLParser
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from kdcube_ai_app.apps.chat.sdk.runtime.workdir_discovery import resolve_output_dir
import kdcube_ai_app.apps.chat.sdk.tools.md_utils as md_utils


# ============================================================================
# CONFIGURATION
# ============================================================================
EMU_PER_INCH = 914400
SLIDE_HEIGHT = Inches(7.5)
SLIDE_WIDTH = Inches(10)
MARGIN = Inches(0.5)
MAX_CONTENT_HEIGHT = Inches(6.0)  # Leave room for title

def page_width_in() -> float:
    return SLIDE_WIDTH.inches - 2 * MARGIN.inches

# Default colors
DEFAULT_COLORS = {
    'text': RGBColor(51, 51, 51),
    'primary': RGBColor(0, 102, 204),
    'subtitle': RGBColor(102, 102, 102),
}


# ============================================================================
# CSS PARSER
# ============================================================================

@dataclass
class StyleInfo:
    color: Optional[RGBColor] = None
    background: Optional[RGBColor] = None
    font_size: Optional[Pt] = None
    line_height: Optional[float] = None  # multiplier, e.g. 1.6
    padding_left: Optional[Inches] = None
    padding_top: Optional[Inches] = None
    padding_right: Optional[Inches] = None
    padding_bottom: Optional[Inches] = None
    border_color: Optional[RGBColor] = None
    border_width: Optional[Pt] = None
    border_left_color: Optional[RGBColor] = None
    border_left_width: Optional[Pt] = None
    bold: bool = False
    italic: bool = False


class CSSParser:
    def __init__(self):
        self.styles: Dict[str, Dict[str, str]] = {}

    def get_style(self, class_name: str) -> StyleInfo:
        """Get computed style for a class."""
        rules = self.styles.get(class_name, {})
        style = StyleInfo()

        if 'color' in rules:
            style.color = self._parse_color(rules['color'])

        if 'background' in rules:
            style.background = self._parse_color(rules['background'])

        if 'font-size' in rules:
            style.font_size = self._parse_font_size(rules['font-size'])

        # Parse border-bottom (for title underline)
        if 'border-bottom' in rules:
            parts = rules['border-bottom'].split()
            if len(parts) >= 3:
                try:
                    # e.g., "4px solid #0066cc"
                    width_str = parts[0].replace('px', '').replace('pt', '')
                    style.border_width = Pt(float(width_str))
                    style.border_color = self._parse_color(parts[2])
                except:
                    pass

        # Parse border-left (for callout boxes)
        if 'border-left' in rules:
            parts = rules['border-left'].split()
            if len(parts) >= 3:
                try:
                    width_str = parts[0].replace('px', '').replace('pt', '')
                    style.border_left_width = Pt(float(width_str))
                    style.border_left_color = self._parse_color(parts[2])
                except:
                    pass

        return style

    def parse(self, css_text: str):
        # Any selector up to '{', rules up to next '}'.
        pattern = r'([^{]+)\{([^}]+)\}'
        for sel, rules_block in re.findall(pattern, css_text):
            # split comma selectors: "th, td"
            for raw in sel.split(','):
                selector = raw.strip()
                # normalize: drop dots/spaces but keep punctuation like ':' for nth-child
                norm = selector.replace('.', '').replace(' ', '')
                rules = {}
                for rule in rules_block.split(';'):
                    rule = rule.strip()
                    if ':' in rule:
                        k, v = rule.split(':', 1)
                        rules[k.strip().lower()] = v.strip()
                # merge if already present
                if norm in self.styles:
                    self.styles[norm].update(rules)
                else:
                    self.styles[norm] = rules

    # NEW: parse inline style="..."
    def parse_inline_rules(self, style_attr: str) -> Dict[str, str]:
        rules = {}
        for piece in style_attr.split(';'):
            piece = piece.strip()
            if ':' in piece:
                k, v = piece.split(':', 1)
                rules[k.strip().lower()] = v.strip()
        return rules

    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        if not color_str:
            return None
        s = color_str.strip().lower()
        if s.startswith('#'):
            s = s[1:]
        if len(s) == 6:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        if len(s) == 3:
            return RGBColor(int(s[0]*2, 16), int(s[1]*2, 16), int(s[2]*2, 16))
        return None  # do not hardcode

    def _parse_font_size(self, size_str: str) -> Optional[Pt]:
        s = size_str.strip().lower()
        if s.endswith('em'):
            return Pt(float(s[:-2]) * 16 * 0.75)  # 1em=16px; 1px=0.75pt
        if s.endswith('px'):
            return Pt(float(s[:-2]) * 0.75)
        if s.endswith('pt'):
            return Pt(float(s[:-2]))
        return None

    def _parse_length_to_inches(self, s: str) -> Optional[Inches]:
        s = s.strip().lower()
        try:
            if s.endswith('px'):  # assume 96dpi for CSS pixel
                return Inches(float(s[:-2]) / 96.0)
            if s.endswith('pt'):
                return Inches(float(s[:-2]) / 72.0)
            if s.endswith('in'):
                return Inches(float(s[:-2]))
        except:
            pass
        return None

    def _inflate_box_shorthand(self, v: str) -> Tuple[Optional[str], Optional[str], Optional[str], Optional[str]]:
        # CSS padding shorthand: t r b l
        parts = v.split()
        if not parts: return (None, None, None, None)
        if len(parts) == 1: parts = parts * 4
        if len(parts) == 2: parts = [parts[0], parts[1], parts[0], parts[1]]
        if len(parts) == 3: parts = [parts[0], parts[1], parts[2], parts[1]]
        return tuple(parts[:4])

    # NEW: merge rules from tag + classes + inline + body with precedence inline>class>tag>body
    def compute_style(self, tag: str, classes: List[str], inline_style: Optional[str]) -> StyleInfo:
        order = []
        if 'body' in self.styles:
            order.append(self.styles['body'])
        if tag in self.styles:
            order.append(self.styles[tag])
        for cls in classes:
            if cls in self.styles:
                order.append(self.styles[cls])
        if inline_style:
            order.append(self.parse_inline_rules(inline_style))

        style = StyleInfo()
        for rules in order:
            if 'color' in rules:
                style.color = self._parse_color(rules['color'])
            if 'background' in rules:
                style.background = self._parse_color(rules['background'])
            if 'background-color' in rules:
                style.background = self._parse_color(rules['background-color'])
            if 'font-size' in rules:
                style.font_size = self._parse_font_size(rules['font-size'])
            if 'line-height' in rules:
                try:
                    style.line_height = float(re.sub('[^0-9.]', '', rules['line-height']))
                except:
                    pass
            if 'padding' in rules:
                t, r, b, l = self._inflate_box_shorthand(rules['padding'])
                style.padding_top = self._parse_length_to_inches(t) if t else style.padding_top
                style.padding_right = self._parse_length_to_inches(r) if r else style.padding_right
                style.padding_bottom = self._parse_length_to_inches(b) if b else style.padding_bottom
                style.padding_left = self._parse_length_to_inches(l) if l else style.padding_left
            for side in ('padding-left','padding-right','padding-top','padding-bottom'):
                if side in rules:
                    val = self._parse_length_to_inches(rules[side])
                    if side == 'padding-left': style.padding_left = val
                    if side == 'padding-right': style.padding_right = val
                    if side == 'padding-top': style.padding_top = val
                    if side == 'padding-bottom': style.padding_bottom = val
            if 'border-bottom' in rules:
                parts = rules['border-bottom'].split()
                if len(parts) >= 3:
                    try:
                        style.border_width = Pt(float(re.sub('[^0-9.]', '', parts[0])))
                        style.border_color = self._parse_color(parts[-1])
                    except:
                        pass
            if 'border-left' in rules:
                parts = rules['border-left'].split()
                if len(parts) >= 3:
                    try:
                        style.border_left_width = Pt(float(re.sub('[^0-9.]', '', parts[0])))
                        style.border_left_color = self._parse_color(parts[-1])
                    except:
                        pass
        return style

def _runs_have_text(runs: List[Dict[str, Any]]) -> bool:
    for r in runs or []:
        if (r.get('text') or '').strip():
            return True
    return False

# ============================================================================
# HTML PARSER
# ============================================================================

class StyledHTMLParser(HTMLParser):
    def __init__(self, css_parser: CSSParser):
        super().__init__()
        self.css = css_parser
        self.body_style = self.css.compute_style('body', [], None)
        self.slides = []
        self.current_section = None
        self.current_element = None
        self.current_list = None
        self.current_list_item = None
        self.in_callout = False  # NEW

        self.format_stack: List[Dict[str, Any]] = []
        self.class_stack: List[List[str]] = []

        self.in_columns = False
        self.columns_container = None    # dict for the whole columns block
        self.current_column_items = None # list for the active column
        self.columns_list = None         # list of columns (each a list of elements)

        self.current_table = None
        self._table_current_row = None
        self._table_in_head = False

        self.div_stack = []
        self.current_column = None

    def _current_format(self) -> Dict[str, Any]:
        return self.format_stack[-1] if self.format_stack else {'bold': False, 'italic': False, 'classes': []}

    def _push_format(self, **kwargs):
        fmt = self._current_format().copy()
        fmt.update(kwargs)
        self.format_stack.append(fmt)

    def _pop_format(self):
        if self.format_stack:
            self.format_stack.pop()
        if self.class_stack:
            self.class_stack.pop()

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        classes = attrs_dict.get('class', '').split()
        inline = attrs_dict.get('style')

        self.class_stack.append(classes)

        if tag in ('strong', 'b'):
            self._push_format(bold=True)
        elif tag in ('em', 'i'):
            self._push_format(italic=True)
        elif tag == 'span':
            self._push_format(classes=classes)

        if self.in_callout and tag in ('h1', 'h2', 'h3'):
            self._push_format(bold=True)
            return

        if tag == 'section':
            self.current_section = {
                'id': attrs_dict.get('id', ''),
                'title': '',
                'subtitle': '',
                'elements': [],
                'title_style': None,
                'subtitle_style': None
            }

        elif tag == 'h1' and self.current_section is not None:
            self.current_element = {
                'type': 'h1',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': self.css.compute_style('h1', classes, inline)
            }

        elif tag == 'h2' and self.current_section is not None:
            self.current_element = {
                'type': 'h2',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': self.css.compute_style('h2', classes, inline)
            }

        elif tag == 'h3' and self.current_section is not None:
            self.current_element = {
                'type': 'h3',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': self.css.compute_style('h3', classes, inline)
            }

        elif tag == 'p' and self.current_section is not None:
            # If we are inside a callout, keep writing into that callout element, do not create a new paragraph element
            if self.in_callout and self.current_element and self.current_element.get('type') == 'callout':
                return
            elem_style = self.css.compute_style('p', classes, inline)
            self.current_element = {
                'type': 'subtitle' if 'subtitle' in classes else 'paragraph',
                'text': '',
                'runs': [],
                'classes': classes,
                'style': elem_style
            }

        elif tag in ('ul', 'ol') and self.current_section is not None:
            self.current_list = {'type': 'list', 'items': []}

        elif tag == 'li' and self.current_list is not None:
            self.current_list_item = {'text': '', 'runs': [], 'style': self.css.compute_style('li', classes, inline)}

        elif tag == 'div' and self.current_section is not None:
            classes = attrs_dict.get('class', '').split()
            inline = attrs_dict.get('style')

            # 1) Two-column container FIRST
            if 'two-column' in classes:
                rules = self.css.styles.get('two-column', {})
                gap_in = 0.0
                if rules and 'gap' in rules:
                    gi = self.css._parse_length_to_inches(rules['gap'])
                    gap_in = gi.inches if gi else 0.0
                self.in_columns = True
                self.columns_list = []
                self.columns_container = {'type': 'columns', 'columns': [], 'gap_in': gap_in}
                self.div_stack.append('two-column')
                return

            # 2) Column inside two-column NEXT (capture style!)
            if self.in_columns and 'column' in classes:
                col_style = self.css.compute_style('div', classes, inline)
                self.current_column = {'type': 'column', 'style': col_style, 'elements': []}
                self.div_stack.append('column')
                return

            # 3) Callouts / highlight AFTER column checks
            elem_style = self.css.compute_style('div', classes, inline)

            # Only promote to callout for explicit callout-ish classes or border-left,
            # not for any background-only container.
            callout_classes = {'highlight', 'highlight-box', 'callout', 'warning', 'success',
                               'phase-box', 'comparison-item'}
            is_callout = any(cls in callout_classes for cls in classes) or bool(elem_style.border_left_color)

            if is_callout:
                self.current_element = {
                    'type': 'callout',
                    'text': '',
                    'runs': [],
                    'classes': classes,
                    'style': elem_style
                }
                self.in_callout = True
                self.div_stack.append('callout')
                return

            # 4) Plain div (no-op container)
            self.div_stack.append('div')

        elif tag == 'table' and self.current_section is not None:
            self.current_table = {'type': 'table', 'header': [], 'rows': [], 'classes': attrs_dict.get('class','').split()}
            self._table_in_head = False

        elif tag == 'thead' and self.current_table is not None:
            self._table_in_head = True

        elif tag == 'tbody' and self.current_table is not None:
            self._table_in_head = False

        elif tag == 'tr' and self.current_table is not None:
            self._table_current_row = []

        elif tag in ('th', 'td') and self.current_table is not None:
            # Create a cell collector
            self.current_element = {'type': 'cell', 'text': '', 'runs': []}

    def handle_endtag(self, tag):
        if tag in ('strong', 'b', 'em', 'i', 'span'):
            self._pop_format()

        if self.in_callout and tag in ('h1', 'h2', 'h3'):
            self._pop_format()
            return

        if tag == 'section' and self.current_section is not None:
            self.slides.append(self.current_section); self.current_section = None

        elif tag == 'h1' and self.current_element:
            self.current_section['title'] = self.current_element['text'].strip()
            self.current_section['title_style'] = self.current_element['style']
            self.current_element = None

        elif tag in ('h2', 'h3') and self.current_element:
            target = (self.current_column['elements'] if (self.in_columns and self.current_column is not None)
                      else self.current_section['elements'])

            target.append(self.current_element)
            self.current_element = None

        elif tag == 'p' and self.current_element:
            if self.in_callout and self.current_element.get('type') == 'callout':
                return
            if self.current_element['type'] == 'subtitle':
                self.current_section['subtitle'] = self.current_element['text'].strip()
                self.current_section['subtitle_style'] = self.current_element['style']
            else:
                target = (self.current_column['elements'] if (self.in_columns and self.current_column is not None)
                          else self.current_section['elements'])

                target.append(self.current_element)
            self.current_element = None

        elif tag in ('ul', 'ol') and self.current_list:
            target = (self.current_column['elements'] if (self.in_columns and self.current_column is not None)
                      else self.current_section['elements'])

            target.append(self.current_list)
            self.current_list = None

        elif tag == 'li' and self.current_list_item:
            self.current_list['items'].append(self.current_list_item); self.current_list_item = None

        elif tag == 'div':
            kind = self.div_stack.pop() if self.div_stack else 'div'

            if kind == 'callout':
                if self.in_columns and self.current_column is not None:
                    self.current_column['elements'].append(self.current_element)
                else:
                    self.current_section['elements'].append(self.current_element)
                self.current_element = None
                self.in_callout = False
                return

            if kind == 'column':
                if self.current_column is not None:
                    self.columns_list.append(self.current_column)
                    self.current_column = None
                return

            if kind == 'two-column':
                # ensure two entries (right may be empty)
                if len(self.columns_list) == 1:
                    self.columns_list.append({'type':'column','style':StyleInfo(), 'elements': []})
                self.columns_container['columns'] = self.columns_list
                self.current_section['elements'].append(self.columns_container)
                self.columns_container = None
                self.columns_list = None
                self.in_columns = False
                return

            return

        elif tag in ('th','td') and self.current_table is not None and self.current_element:
            # finalize cell
            self._table_current_row.append(self.current_element)
            self.current_element = None

        elif tag == 'tr' and self.current_table is not None and self._table_current_row is not None:
            cells = self._table_current_row
            texts = []
            for c in cells:
                texts.append({'text': c['text'], 'runs': c['runs']})
            if self._table_in_head and not self.current_table['header']:
                self.current_table['header'] = texts
            else:
                self.current_table['rows'].append(texts)
            self._table_current_row = None

        elif tag == 'table' and self.current_table is not None:
            target = (self.current_column['elements'] if (self.in_columns and self.current_column is not None)
                      else self.current_section['elements'])
            target.append(self.current_table)
            self.current_table = None

    def handle_data(self, data):
        if not data.strip():
            return
        fmt = self._current_format()
        # Inherit color: inline class run color > element style color > body color
        color = None
        for cls in fmt.get('classes', []):
            s = self.css.get_style(cls)  # still okay for span classes
            if s.color:
                color = s.color; break
        if color is None and self.current_element:
            elem_style = self.current_element.get('style')
            if elem_style and elem_style.color:
                color = elem_style.color
        if color is None and self.body_style.color:
            color = self.body_style.color

        run = {'text': data, 'bold': fmt.get('bold', False), 'italic': fmt.get('italic', False), 'color': color}

        if self.current_list_item is not None:
            self.current_list_item['text'] += data
            self.current_list_item['runs'].append(run)
        elif self.current_element is not None:
            self.current_element['text'] += data
            if 'runs' in self.current_element:
                self.current_element['runs'].append(run)


# ============================================================================
# RENDERER WITH OVERFLOW PROTECTION
# ============================================================================
def _estimate_elements_height(elements: List[Dict], base_font_size_pt: float, page_w_in: float) -> float:
    total = 0.0
    for elem in elements:
        t = elem.get('type')
        style: StyleInfo = elem.get('style') or StyleInfo()
        if t in ('h2','h3'):
            fs = (style.font_size.pt if style.font_size else (28 if t=='h2' else 22))
            h = _estimate_runs_height_in_inches(elem.get('runs', []), fs, page_w_in, (style.line_height or 1.2))
            total += h.inches + 0.1
        elif t == 'paragraph':
            fs = (style.font_size.pt if style.font_size else 18)
            h = _estimate_runs_height_in_inches(elem.get('runs', []), fs, page_w_in, (style.line_height or 1.3))
            total += h.inches + 0.06
        elif t == 'list':
            fs = base_font_size_pt
            w = page_w_in - 0.25
            for it in elem.get('items', []):
                h = _estimate_runs_height_in_inches(it.get('runs', []), fs, w, 1.3)
                total += h.inches + 0.06
            total += 0.06
        elif t == 'callout':
            fs = (style.font_size.pt if style.font_size else 16)
            pad_l = (style.padding_left or Inches(0.2)).inches
            pad_r = (style.padding_right or Inches(0.1)).inches
            pad_t = (style.padding_top or Inches(0.1)).inches
            pad_b = (style.padding_bottom or Inches(0.1)).inches
            w = page_w_in - pad_l - pad_r
            h = _estimate_runs_height_in_inches(elem.get('runs', []), fs, w, (style.line_height or 1.3))
            total += max(h.inches + pad_t + pad_b, 0.6) + 0.12
        elif t == 'table':
            cols = len(elem.get('header', [])) or (len(elem.get('rows', [])[0]) if elem.get('rows') else 0)
            if cols == 0:
                continue
            fs = 14
            header_h = (fs * 1.2) / 72.0
            row_h = (fs * 1.2) / 72.0
            total += header_h + row_h * len(elem.get('rows', [])) + 0.1
    return total


def _avg_char_width_pt(font_pt: float, bold: bool) -> float:
    # rough average, works well enough for layout
    base = font_pt * 0.52
    return base * (1.08 if bold else 1.0)

def _estimate_runs_height_in_inches(runs: List[Dict[str,Any]], font_pt: float, width_in_inches: float, line_height_mult: float = 1.2) -> Inches:
    if not runs:
        return Inches((font_pt * line_height_mult) / 72.0)
    width_pts = width_in_inches * 72.0
    text = ''.join(r['text'] for r in runs)
    # estimate chars per line using weighted average width across runs
    if not text.strip():
        return Inches((font_pt * line_height_mult) / 72.0)

    # conservative: use the smallest chars-per-line across runs
    cpl_candidates = []
    for r in runs:
        fw = _avg_char_width_pt(font_pt, r.get('bold', False))
        if fw <= 0: continue
        cpl_candidates.append(int(max(1, width_pts / fw)))
    cpl = min(cpl_candidates) if cpl_candidates else int(max(1, width_pts / _avg_char_width_pt(font_pt, False)))

    # rough wrap count
    lines = 0
    for paragraph in text.split('\n'):
        length = max(1, len(paragraph))
        lines += int((length + cpl - 1) // cpl)
    if lines < 1: lines = 1
    height_pts = lines * (font_pt * line_height_mult)
    return Inches(height_pts / 72.0)


def estimate_content_height(elements: List[Dict], base_font_size: Pt) -> float:
    total_in = 0.0
    page_w_in = SLIDE_WIDTH.inches - 2 * MARGIN.inches  # <-- float inches

    for elem in elements:
        if not elem:
            continue
        t = elem.get('type')
        style: StyleInfo = elem.get('style') or StyleInfo()

        if t in ('h2', 'h3'):
            fs = (style.font_size or (Pt(28) if t == 'h2' else Pt(22))).pt
            h = _estimate_runs_height_in_inches(
                elem.get('runs', []),
                fs,
                page_w_in,
                (style.line_height or 1.2),
            )
            total_in += h.inches + 0.1

        elif t == 'list':
            fs = base_font_size.pt
            bullets_w_in = page_w_in - 0.25  # 0.25" indent on the left
            for it in elem.get('items', []):
                runs = it.get('runs', [])
                h = _estimate_runs_height_in_inches(runs, fs, bullets_w_in, 1.3)
                total_in += h.inches + 0.06
            total_in += 0.1

        elif t == 'callout':
            fs = (style.font_size or Pt(16)).pt
            pad_l_in = (style.padding_left or Inches(0.2)).inches
            pad_r_in = (style.padding_right or Inches(0.1)).inches
            pad_t_in = (style.padding_top or Inches(0.1)).inches
            pad_b_in = (style.padding_bottom or Inches(0.1)).inches

            content_w_in = page_w_in - pad_l_in - pad_r_in
            h = _estimate_runs_height_in_inches(
                elem.get('runs', []),
                fs,
                content_w_in,
                (style.line_height or 1.3),
            )
            total_in += max(h.inches + pad_t_in + pad_b_in, 0.6) + 0.15
        elif t == 'columns':
            gap_in = elem.get('gap_in', 0.0)
            col_w_in = (page_w_in - gap_in) / 2.0
            cols = elem.get('columns', [])
            def col_height(col):
                if not col: return 0.0
                cs: StyleInfo = col.get('style') or StyleInfo()
                pad_l = (cs.padding_left or Inches(0.0)).inches
                pad_r = (cs.padding_right or Inches(0.0)).inches
                pad_t = (cs.padding_top or Inches(0.0)).inches
                pad_b = (cs.padding_bottom or Inches(0.0)).inches
                inner_w = max(0.1, col_w_in - pad_l - pad_r)
                content_h = _estimate_elements_height(col.get('elements', []), base_font_size.pt, inner_w)
                return content_h + pad_t + pad_b
            left_h = col_height(cols[0]) if len(cols) > 0 else 0.0
            right_h = col_height(cols[1]) if len(cols) > 1 else 0.0
            total_in += max(left_h, right_h) + 0.12

        elif t == 'table':
            cols = len(elem.get('header', [])) or (len(elem.get('rows', [])[0]) if elem.get('rows') else 0)
            if cols:
                fs = 14
                header_h = (fs * 1.2) / 72.0
                row_h = (fs * 1.2) / 72.0
                total_in += header_h + row_h * len(elem.get('rows', [])) + 0.1
    return total_in


# def render_slide(
#         prs: Presentation,
#         slide_data: Dict[str, Any],
#         css_parser,
#         sources_map: Dict[int, Dict[str, str]],
#         resolve_citations: bool
# ):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#
#     title_text = slide_data.get('title', '')
#     title_style: StyleInfo = slide_data.get('title_style') or StyleInfo()
#     subtitle_text = slide_data.get('subtitle', '')
#     subtitle_style: StyleInfo = slide_data.get('subtitle_style') or StyleInfo()
#     elements = slide_data.get('elements', [])
#
#     page_w_in = SLIDE_WIDTH.inches - 2 * MARGIN.inches  # float inches
#
#     estimated_height = estimate_content_height(elements, Pt(18))
#     scale_factor = 1.0
#     if estimated_height > MAX_CONTENT_HEIGHT.inches:
#         scale_factor = max(0.7, MAX_CONTENT_HEIGHT.inches / estimated_height)
#
#     y = MARGIN
#
#     # TITLE
#     title_h = Inches(0.8 * scale_factor)
#     title_box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, title_h)
#     tf = title_box.text_frame
#     tf.word_wrap = True
#     p = tf.paragraphs[0]
#     r = p.add_run(); r.text = title_text
#     r.font.size = title_style.font_size or Pt(36 * scale_factor)
#     r.font.bold = True
#     if title_style.color:
#         r.font.color.rgb = title_style.color
#     y += title_h
#
#     # underline from CSS
#     if title_style.border_color and title_style.border_width:
#         line = slide.shapes.add_shape(
#             MSO_SHAPE.RECTANGLE,
#             MARGIN, y,
#             SLIDE_WIDTH - MARGIN * 2, title_style.border_width
#         )
#         line.fill.solid()
#         line.fill.fore_color.rgb = title_style.border_color
#         line.line.fill.background()
#         y += Inches(0.12 * scale_factor)
#
#     # SUBTITLE
#     if subtitle_text:
#         sub_h = Inches(0.45 * scale_factor)
#         subtitle_box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, sub_h)
#         tf = subtitle_box.text_frame
#         p = tf.paragraphs[0]
#         # Use sources-aware emitter for any links/citations in subtitle
#         _emit_text_with_links_and_citations(
#             p,
#             text=subtitle_text,
#             default_size_pt=(subtitle_style.font_size.pt if subtitle_style.font_size else 18 * scale_factor),
#             bold=False,
#             italic=bool(subtitle_style.italic),
#             color=(subtitle_style.color or DEFAULT_COLORS.get('subtitle')),
#             sources_map=sources_map,
#             resolve_citations=resolve_citations
#         )
#         y += sub_h + Inches(0.12 * scale_factor)
#
#     # CONTENT
#     base_font_pt = 18 * scale_factor
#
#     for elem in elements:
#         if not elem:
#             continue
#         elem_type = elem.get('type')
#         style: StyleInfo = elem.get('style') or StyleInfo()
#         runs = elem.get('runs', [])
#
#         if elem_type in ('h2', 'h3'):
#             fs = (style.font_size.pt if style.font_size else (28 if elem_type == 'h2' else 22)) * scale_factor
#             box_h = _estimate_runs_height_in_inches(runs, fs, page_w_in, (style.line_height or 1.2))
#             box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, box_h)
#             tf = box.text_frame; tf.word_wrap = True
#             p = tf.paragraphs[0]
#             # Force bold for headings, but still resolve links/citations inside
#             for run_data in runs or [{'text': ''}]:
#                 _emit_text_with_links_and_citations(
#                     p,
#                     text=run_data.get('text', ''),
#                     default_size_pt=fs,
#                     bold=True,
#                     italic=bool(run_data.get('italic', False)),
#                     color=(run_data.get('color') or style.color or DEFAULT_COLORS.get('text')),
#                     sources_map=sources_map,
#                     resolve_citations=resolve_citations
#                 )
#             y += box_h + Inches(0.08 * scale_factor)
#
#         elif elem_type == 'list':
#             # keep separate measurement width vs drawing width
#             w_len = SLIDE_WIDTH - MARGIN * 2 - Inches(0.25)     # Length/EMU for drawing
#             w_in  = page_w_in - 0.25                            # float inches for measurement
#             for item in elem.get('items', []):
#                 runs_i = item.get('runs', [])
#                 fs = base_font_pt
#                 lh = 1.3
#                 h = _estimate_runs_height_in_inches(runs_i, fs, w_in, lh)
#                 box = slide.shapes.add_textbox(MARGIN + Inches(0.25), y, w_len, h)
#                 tf = box.text_frame; tf.word_wrap = True
#                 p = tf.paragraphs[0]; p.level = 0
#                 _emit_runs_with_sources(
#                     p,
#                     runs_i,
#                     default_size_pt=fs,
#                     fallback_color=DEFAULT_COLORS.get('text'),
#                     sources_map=sources_map,
#                     resolve_citations=resolve_citations
#                 )
#                 y += h + Inches(0.06 * scale_factor)
#             y += Inches(0.06 * scale_factor)
#
#         elif elem_type == 'callout':
#             # Skip ghost callouts with no visible text
#             if not _runs_have_text(runs):
#                 continue
#             fs = (style.font_size.pt if style.font_size else 16) * scale_factor
#             pad_l = style.padding_left or Inches(0.2)
#             pad_r = style.padding_right or Inches(0.1)
#             pad_t = style.padding_top or Inches(0.1)
#             pad_b = style.padding_bottom or Inches(0.1)
#
#             content_w_in = page_w_in - pad_l.inches - pad_r.inches
#             content_w_len = SLIDE_WIDTH - MARGIN * 2 - pad_l - pad_r  # for drawing
#
#             content_h = _estimate_runs_height_in_inches(runs, fs, content_w_in, (style.line_height or 1.3))
#             box_h_in = max(content_h.inches + pad_t.inches + pad_b.inches, 0.6)
#             box_h = Inches(box_h_in)
#
#             if style.background:
#                 bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, SLIDE_WIDTH - MARGIN * 2, box_h)
#                 bg.fill.solid(); bg.fill.fore_color.rgb = style.background
#                 bg.line.fill.background()
#
#             if style.border_left_color and style.border_left_width:
#                 border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, style.border_left_width, box_h)
#                 border.fill.solid(); border.fill.fore_color.rgb = style.border_left_color
#                 border.line.fill.background()
#
#             text_box = slide.shapes.add_textbox(MARGIN + pad_l, y + pad_t, content_w_len, box_h - pad_t - pad_b)
#             tf = text_box.text_frame; tf.word_wrap = True
#             p = tf.paragraphs[0]
#             _emit_runs_with_sources(
#                 p,
#                 runs,
#                 default_size_pt=fs,
#                 fallback_color=(style.color or DEFAULT_COLORS.get('text')),
#                 sources_map=sources_map,
#                 resolve_citations=resolve_citations
#             )
#             y += box_h + Inches(0.12 * scale_factor)
#
#         elif elem_type == 'columns':
#             gap_in = elem.get('gap_in', 0.0)
#             gap_len = Inches(gap_in)
#             col_w_in = (page_w_in - gap_in) / 2.0
#             col_w_len = Inches(col_w_in)
#             left_x = MARGIN
#             right_x = MARGIN + col_w_len + gap_len
#             cols = elem.get('columns', [])
#
#             def compute_col_dims(col):
#                 cs: StyleInfo = (col.get('style') if col else None) or StyleInfo()
#                 pad_l = cs.padding_left or Inches(0.0)
#                 pad_r = cs.padding_right or Inches(0.0)
#                 pad_t = cs.padding_top or Inches(0.0)
#                 pad_b = cs.padding_bottom or Inches(0.0)
#                 inner_w_in = max(0.1, col_w_in - pad_l.inches - pad_r.inches)
#                 content_h_in = _estimate_elements_height(col.get('elements', []), 18 * scale_factor, inner_w_in) if col else 0.0
#                 total_h_in = content_h_in + pad_t.inches + pad_b.inches
#                 return cs, pad_l, pad_r, pad_t, pad_b, inner_w_in, Inches(inner_w_in), Inches(total_h_in)
#
#             # precompute both columns
#             colL = cols[0] if len(cols) > 0 else {'style':StyleInfo(), 'elements':[]}
#             colR = cols[1] if len(cols) > 1 else {'style':StyleInfo(), 'elements':[]}
#             csL, padL_l, padL_r, padL_t, padL_b, innerL_in, innerL_len, colH_len_L = compute_col_dims(colL)
#             csR, padR_l, padR_r, padR_t, padR_b, innerR_in, innerR_len, colH_len_R = compute_col_dims(colR)
#
#             # draw backgrounds first (only if there is some content height)
#             min_draw_in = 0.05
#             if csL.background and colH_len_L.inches > min_draw_in:
#                 bgL = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y, col_w_len, colH_len_L)
#                 bgL.fill.solid(); bgL.fill.fore_color.rgb = csL.background
#                 bgL.line.fill.background()
#             if csR.background and colH_len_R.inches > min_draw_in:
#                 bgR = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y, col_w_len, colH_len_R)
#                 bgR.fill.solid(); bgR.fill.fore_color.rgb = csR.background
#                 bgR.line.fill.background()
#
#             # optional left borders
#             if csL.border_left_color and csL.border_left_width:
#                 bL = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y, csL.border_left_width, colH_len_L)
#                 bL.fill.solid(); bL.fill.fore_color.rgb = csL.border_left_color; bL.line.fill.background()
#             if csR.border_left_color and csR.border_left_width:
#                 bR = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y, csR.border_left_width, colH_len_R)
#                 bR.fill.solid(); bR.fill.fore_color.rgb = csR.border_left_color; bR.line.fill.background()
#
#             def render_into(x_left_len, y_top_len, avail_w_len, avail_w_in, elements_in_col):
#                 y_col = y_top_len
#                 base_pt = 18 * scale_factor
#                 for b in elements_in_col:
#                     bt = b.get('type'); bs: StyleInfo = b.get('style') or StyleInfo()
#                     if bt in ('h2','h3'):
#                         fs = (bs.font_size.pt if bs.font_size else (28 if bt=='h2' else 22)) * scale_factor
#                         h = _estimate_runs_height_in_inches(b.get('runs', []), fs, avail_w_in, (bs.line_height or 1.2))
#                         box = slide.shapes.add_textbox(x_left_len, y_col, avail_w_len, h)
#                         tf = box.text_frame; tf.word_wrap = True
#                         p = tf.paragraphs[0]
#                         # Force bold for headings inside columns
#                         for rd in b.get('runs', []) or [{'text': ''}]:
#                             _emit_text_with_links_and_citations(
#                                 p,
#                                 text=rd.get('text', ''),
#                                 default_size_pt=fs,
#                                 bold=True,
#                                 italic=bool(rd.get('italic', False)),
#                                 color=(rd.get('color') or bs.color or DEFAULT_COLORS.get('text')),
#                                 sources_map=sources_map,
#                                 resolve_citations=resolve_citations
#                             )
#                         y_col += h + Inches(0.06 * scale_factor)
#
#                     elif bt == 'paragraph':
#                         fs = (bs.font_size.pt if bs.font_size else 18) * scale_factor
#                         h = _estimate_runs_height_in_inches(b.get('runs', []), fs, avail_w_in, (bs.line_height or 1.3))
#                         box = slide.shapes.add_textbox(x_left_len, y_col, avail_w_len, h)
#                         tf = box.text_frame; tf.word_wrap = True
#                         p = tf.paragraphs[0]
#                         _emit_runs_with_sources(
#                             p,
#                             b.get('runs', []),
#                             default_size_pt=fs,
#                             fallback_color=(bs.color or DEFAULT_COLORS.get('text')),
#                             sources_map=sources_map,
#                             resolve_citations=resolve_citations
#                         )
#                         y_col += h + Inches(0.04 * scale_factor)
#
#                     elif bt == 'list':
#                         fs = base_pt
#                         w_len = avail_w_len - Inches(0.25)
#                         w_in  = max(0.1, avail_w_in - 0.25)
#                         for it in b.get('items', []):
#                             h = _estimate_runs_height_in_inches(it.get('runs', []), fs, w_in, 1.3)
#                             box = slide.shapes.add_textbox(x_left_len + Inches(0.25), y_col, w_len, h)
#                             tf = box.text_frame; tf.word_wrap = True
#                             p = tf.paragraphs[0]; p.level = 0
#                             _emit_runs_with_sources(
#                                 p,
#                                 it.get('runs', []),
#                                 default_size_pt=fs,
#                                 fallback_color=DEFAULT_COLORS.get('text'),
#                                 sources_map=sources_map,
#                                 resolve_citations=resolve_citations
#                             )
#                             y_col += h + Inches(0.04 * scale_factor)
#                         y_col += Inches(0.02 * scale_factor)
#
#                     elif bt == 'callout':
#                         fs = (bs.font_size.pt if bs.font_size else 16) * scale_factor
#                         pad_l = bs.padding_left or Inches(0.2)
#                         pad_r = bs.padding_right or Inches(0.1)
#                         pad_t = bs.padding_top or Inches(0.1)
#                         pad_b = bs.padding_bottom or Inches(0.1)
#                         inner_w_in  = max(0.1, avail_w_in - pad_l.inches - pad_r.inches)
#                         inner_w_len = avail_w_len - pad_l - pad_r
#                         content_h = _estimate_runs_height_in_inches(b.get('runs', []), fs, inner_w_in, (bs.line_height or 1.3))
#                         box_h = Inches(max(content_h.inches + pad_t.inches + pad_b.inches, 0.6))
#
#                         if bs.background:
#                             bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_left_len, y_col, avail_w_len, box_h)
#                             bg.fill.solid(); bg.fill.fore_color.rgb = bs.background; bg.line.fill.background()
#                         if bs.border_left_color and bs.border_left_width:
#                             border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_left_len, y_col, bs.border_left_width, box_h)
#                             border.fill.solid(); border.fill.fore_color.rgb = bs.border_left_color; border.line.fill.background()
#
#                         tb = slide.shapes.add_textbox(x_left_len + pad_l, y_col + pad_t, inner_w_len, box_h - pad_t - pad_b)
#                         tf = tb.text_frame; tf.word_wrap = True
#                         p = tf.paragraphs[0]
#                         _emit_runs_with_sources(
#                             p,
#                             b.get('runs', []),
#                             default_size_pt=fs,
#                             fallback_color=(bs.color or DEFAULT_COLORS.get('text')),
#                             sources_map=sources_map,
#                             resolve_citations=resolve_citations
#                         )
#                         y_col += box_h + Inches(0.08 * scale_factor)
#
#                     elif bt == 'table':
#                         header = b.get('header', [])
#                         rows = b.get('rows', [])
#
#                         # promote first row to header if needed
#                         if not header and rows:
#                             header, rows = rows[0], rows[1:]
#
#                         if not header:
#                             continue
#
#                         ncols = len(header)
#
#                         # heights tuned a bit smaller for columns
#                         row_h = Inches(0.30 * scale_factor)
#                         header_h = Inches(0.35 * scale_factor)
#
#                         # colors from CSS
#                         th_rules = css_parser.styles.get('th', {})
#                         stripe_rules = css_parser.styles.get('tr:nth-child(even)', {})
#                         parse_color = CSSParser()._parse_color
#                         header_bg = parse_color(th_rules.get('background-color')) if th_rules else None
#                         header_fg = parse_color(th_rules.get('color')) if th_rules else None
#                         stripe_bg = parse_color(stripe_rules.get('background-color')) if stripe_rules else None
#
#                         tbl = slide.shapes.add_table(1 + len(rows), ncols,
#                                                      x_left_len, y_col,
#                                                      avail_w_len, header_h + row_h*max(1, len(rows))).table
#
#                         # header
#                         for c, cell in enumerate(header):
#                             tc = tbl.cell(0, c)
#                             tf = tc.text_frame; tf.clear()
#                             p = tf.paragraphs[0]
#                             _emit_runs_with_sources(
#                                 p,
#                                 cell.get('runs', []),
#                                 default_size_pt=14 * scale_factor,
#                                 fallback_color=(header_fg or DEFAULT_COLORS.get('text')),
#                                 sources_map=sources_map,
#                                 resolve_citations=resolve_citations
#                             )
#                             for rr in p.runs: rr.font.bold = True
#                             if header_bg:
#                                 tc.fill.solid(); tc.fill.fore_color.rgb = header_bg
#
#                         # body
#                         for r, row in enumerate(rows, start=1):
#                             for c, cell in enumerate(row):
#                                 tc = tbl.cell(r, c)
#                                 tf = tc.text_frame; tf.clear()
#                                 p = tf.paragraphs[0]
#                                 _emit_runs_with_sources(
#                                     p,
#                                     cell.get('runs', []),
#                                     default_size_pt=14 * scale_factor,
#                                     fallback_color=DEFAULT_COLORS.get('text'),
#                                     sources_map=sources_map,
#                                     resolve_citations=resolve_citations
#                                 )
#                             if stripe_bg and (r % 2 == 0):
#                                 for c in range(ncols):
#                                     tc = tbl.cell(r, c); tc.fill.solid(); tc.fill.fore_color.rgb = stripe_bg
#
#                         y_col += header_h + row_h*max(1, len(rows)) + Inches(0.06 * scale_factor)
#
#                 return y_col
#
#             # content frames top-lefts (apply column paddings)
#             yL_start = y + padL_t
#             yR_start = y + padR_t
#             xL_inset = left_x + padL_l
#             xR_inset = right_x + padR_l
#             wL_inner = col_w_len - padL_l - padL_r
#             wR_inner = col_w_len - padR_l - padR_r
#
#             y_left_end  = render_into(xL_inset, yL_start, wL_inner, innerL_in, colL.get('elements', []))
#             y_right_end = render_into(xR_inset, yR_start, wR_inner, innerR_in, colR.get('elements', []))
#
#             # advance y by the taller of the two columns
#             y += Inches(max(colH_len_L.inches, colH_len_R.inches)) + Inches(0.04 * scale_factor)
#
#         elif elem_type == 'table':
#             header = elem.get('header', [])
#             rows = elem.get('rows', [])
#
#             # promote first row to header if no <thead> was present
#             if not header and rows:
#                 header, rows = rows[0], rows[1:]
#
#             if header:
#                 ncols = len(header)
#
#                 # width + heights
#                 tbl_w_len = SLIDE_WIDTH - MARGIN*2
#                 row_h = Inches(0.35 * scale_factor)
#                 header_h = Inches(0.4 * scale_factor)
#
#                 # read CSS colors directly from the parser
#                 th_rules = css_parser.styles.get('th', {})
#                 stripe_rules = css_parser.styles.get('tr:nth-child(even)', {})
#                 parse_color = CSSParser()._parse_color
#                 header_bg = parse_color(th_rules.get('background-color')) if th_rules else None
#                 header_fg = parse_color(th_rules.get('color')) if th_rules else None
#                 stripe_bg = parse_color(stripe_rules.get('background-color')) if stripe_rules else None
#
#                 tbl = slide.shapes.add_table(1 + len(rows), ncols, MARGIN, y,
#                                              tbl_w_len, header_h + row_h*max(1, len(rows))).table
#
#                 # header cells
#                 for c, cell in enumerate(header):
#                     tc = tbl.cell(0, c)
#                     tf = tc.text_frame; tf.clear()
#                     p = tf.paragraphs[0]
#                     _emit_runs_with_sources(
#                         p,
#                         cell.get('runs', []),
#                         default_size_pt=14 * scale_factor,
#                         fallback_color=(header_fg or DEFAULT_COLORS.get('text')),
#                         sources_map=sources_map,
#                         resolve_citations=resolve_citations
#                     )
#                     for rr in p.runs: rr.font.bold = True
#                     if header_bg:
#                         tc.fill.solid(); tc.fill.fore_color.rgb = header_bg
#
#                 # body rows
#                 for r, row in enumerate(rows, start=1):
#                     for c, cell in enumerate(row):
#                         tc = tbl.cell(r, c)
#                         tf = tc.text_frame; tf.clear()
#                         p = tf.paragraphs[0]
#                         _emit_runs_with_sources(
#                             p,
#                             cell.get('runs', []),
#                             default_size_pt=14 * scale_factor,
#                             fallback_color=DEFAULT_COLORS.get('text'),
#                             sources_map=sources_map,
#                             resolve_citations=resolve_citations
#                         )
#                     if stripe_bg and (r % 2 == 0):
#                         for c in range(ncols):
#                             tc = tbl.cell(r, c); tc.fill.solid(); tc.fill.fore_color.rgb = stripe_bg
#
#                 y += header_h + row_h*max(1, len(rows)) + Inches(0.12 * scale_factor)
#
#         elif elem_type == 'paragraph':
#             # (Top-level paragraph — previously missing; safe to handle)
#             fs = (style.font_size.pt if style.font_size else 18) * scale_factor
#             h = _estimate_runs_height_in_inches(runs, fs, page_w_in, (style.line_height or 1.3))
#             box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - 2*MARGIN, h)
#             tf = box.text_frame; tf.word_wrap = True
#             p = tf.paragraphs[0]
#             _emit_runs_with_sources(
#                 p,
#                 runs,
#                 default_size_pt=fs,
#                 fallback_color=(style.color or DEFAULT_COLORS.get('text')),
#                 sources_map=sources_map,
#                 resolve_citations=resolve_citations
#             )
#             y += h + Inches(0.08 * scale_factor)
def render_slide(
        prs: Presentation,
        slide_data: Dict[str, Any],
        css_parser,
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_text = slide_data.get('title', '')
    title_style: StyleInfo = slide_data.get('title_style') or StyleInfo()
    subtitle_text = slide_data.get('subtitle', '')
    subtitle_style: StyleInfo = slide_data.get('subtitle_style') or StyleInfo()
    elements = slide_data.get('elements', [])

    page_w_in = SLIDE_WIDTH.inches - 2 * MARGIN.inches

    # Scale based on content (body only), same as before
    estimated_height = estimate_content_height(elements, Pt(18))
    scale_factor = 1.0
    if estimated_height > MAX_CONTENT_HEIGHT.inches:
        scale_factor = max(0.7, MAX_CONTENT_HEIGHT.inches / estimated_height)

    y = MARGIN

    # === TITLE (measure, then draw) ==========================================
    title_fs_pt = (title_style.font_size.pt if title_style.font_size else 36 * scale_factor)
    title_runs_for_measure = [{'text': title_text or ''}]
    title_h = _estimate_runs_height_in_inches(
        title_runs_for_measure,
        title_fs_pt,
        page_w_in,
        (title_style.line_height or 1.15),
    )
    # keep a sensible minimum but let large titles expand
    if title_h.inches < 0.6 * scale_factor:
        title_h = Inches(0.6 * scale_factor)

    title_box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, title_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title_text
    r.font.size = Pt(title_fs_pt)
    r.font.bold = True
    if title_style.color:
        r.font.color.rgb = title_style.color
    y += title_h

    # Underline from CSS — place it **after** measured title; then add a gap
    if title_style.border_color and title_style.border_width:
        line_h = title_style.border_width
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN, y,
            SLIDE_WIDTH - MARGIN * 2, line_h
        ).line.fill.background()  # line has no stroke; the fill is enough
        # color fill
        underline = slide.shapes[-1]
        underline.fill.solid()
        underline.fill.fore_color.rgb = title_style.border_color
        y += Inches(0.15 * scale_factor)  # a touch more air than before

    # === SUBTITLE (measure, then draw) =======================================
    if subtitle_text:
        sub_fs_pt = (subtitle_style.font_size.pt if subtitle_style.font_size else 18 * scale_factor)
        sub_h = _estimate_runs_height_in_inches(
            [{'text': subtitle_text}],
            sub_fs_pt,
            page_w_in,
            (subtitle_style.line_height or 1.2),
        )
        if sub_h.inches < 0.35 * scale_factor:
            sub_h = Inches(0.35 * scale_factor)

        subtitle_box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, sub_h)
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        _emit_text_with_links_and_citations(
            p,
            text=subtitle_text,
            default_size_pt=sub_fs_pt,
            bold=False,
            italic=bool(subtitle_style.italic),
            color=(subtitle_style.color or DEFAULT_COLORS.get('subtitle')),
            sources_map=sources_map,
            resolve_citations=resolve_citations
        )
        y += sub_h + Inches(0.12 * scale_factor)

    # === BODY ================================================================
    base_font_pt = 18 * scale_factor

    for elem in elements:
        if not elem:
            continue
        elem_type = elem.get('type')
        style: StyleInfo = elem.get('style') or StyleInfo()
        runs = elem.get('runs', [])

        if elem_type in ('h2', 'h3'):
            fs = (style.font_size.pt if style.font_size else (28 if elem_type == 'h2' else 22)) * scale_factor
            box_h = _estimate_runs_height_in_inches(runs, fs, page_w_in, (style.line_height or 1.2))
            box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - MARGIN * 2, box_h)
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            for run_data in runs or [{'text': ''}]:
                _emit_text_with_links_and_citations(
                    p,
                    text=run_data.get('text', ''),
                    default_size_pt=fs,
                    bold=True,
                    italic=bool(run_data.get('italic', False)),
                    color=(run_data.get('color') or style.color or DEFAULT_COLORS.get('text')),
                    sources_map=sources_map,
                    resolve_citations=resolve_citations
                )
            y += box_h + Inches(0.08 * scale_factor)

        elif elem_type == 'list':
            w_len = SLIDE_WIDTH - MARGIN * 2 - Inches(0.25)
            w_in  = page_w_in - 0.25
            for item in elem.get('items', []):
                runs_i = item.get('runs', [])
                fs = base_font_pt
                lh = 1.3
                h = _estimate_runs_height_in_inches(runs_i, fs, w_in, lh)
                box = slide.shapes.add_textbox(MARGIN + Inches(0.25), y, w_len, h)
                tf = box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.level = 0
                _emit_runs_with_sources(
                    p,
                    runs_i,
                    default_size_pt=fs,
                    fallback_color=DEFAULT_COLORS.get('text'),
                    sources_map=sources_map,
                    resolve_citations=resolve_citations
                )
                y += h + Inches(0.06 * scale_factor)
            y += Inches(0.06 * scale_factor)

        elif elem_type == 'callout':
            if not _runs_have_text(runs):
                continue
            fs = (style.font_size.pt if style.font_size else 16) * scale_factor
            pad_l = style.padding_left or Inches(0.2)
            pad_r = style.padding_right or Inches(0.1)
            pad_t = style.padding_top or Inches(0.1)
            pad_b = style.padding_bottom or Inches(0.1)

            content_w_in = page_w_in - pad_l.inches - pad_r.inches
            content_w_len = SLIDE_WIDTH - MARGIN * 2 - pad_l - pad_r

            content_h = _estimate_runs_height_in_inches(runs, fs, content_w_in, (style.line_height or 1.3))
            box_h_in = max(content_h.inches + pad_t.inches + pad_b.inches, 0.6)
            box_h = Inches(box_h_in)

            if style.background:
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, SLIDE_WIDTH - MARGIN * 2, box_h)
                bg.fill.solid(); bg.fill.fore_color.rgb = style.background
                bg.line.fill.background()

            if style.border_left_color and style.border_left_width:
                border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, style.border_left_width, box_h)
                border.fill.solid(); border.fill.fore_color.rgb = style.border_left_color
                border.line.fill.background()

            text_box = slide.shapes.add_textbox(MARGIN + pad_l, y + pad_t, content_w_len, box_h - pad_t - pad_b)
            tf = text_box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            _emit_runs_with_sources(
                p,
                runs,
                default_size_pt=fs,
                fallback_color=(style.color or DEFAULT_COLORS.get('text')),
                sources_map=sources_map,
                resolve_citations=resolve_citations
            )
            y += box_h + Inches(0.12 * scale_factor)

        elif elem_type == 'columns':
            gap_in = elem.get('gap_in', 0.0)
            gap_len = Inches(gap_in)
            col_w_in = (page_w_in - gap_in) / 2.0
            col_w_len = Inches(col_w_in)
            left_x = MARGIN
            right_x = MARGIN + col_w_len + gap_len
            cols = elem.get('columns', [])

            def compute_col_dims(col):
                cs: StyleInfo = (col.get('style') if col else None) or StyleInfo()
                pad_l = cs.padding_left or Inches(0.0)
                pad_r = cs.padding_right or Inches(0.0)
                pad_t = cs.padding_top or Inches(0.0)
                pad_b = cs.padding_bottom or Inches(0.0)
                inner_w_in = max(0.1, col_w_in - pad_l.inches - pad_r.inches)
                content_h_in = _estimate_elements_height(col.get('elements', []), 18 * scale_factor, inner_w_in) if col else 0.0
                total_h_in = content_h_in + pad_t.inches + pad_b.inches
                return cs, pad_l, pad_r, pad_t, pad_b, inner_w_in, Inches(inner_w_in), Inches(total_h_in)

            colL = cols[0] if len(cols) > 0 else {'style':StyleInfo(), 'elements':[]}
            colR = cols[1] if len(cols) > 1 else {'style':StyleInfo(), 'elements':[]}
            csL, padL_l, padL_r, padL_t, padL_b, innerL_in, innerL_len, colH_len_L = compute_col_dims(colL)
            csR, padR_l, padR_r, padR_t, padR_b, innerR_in, innerR_len, colH_len_R = compute_col_dims(colR)

            min_draw_in = 0.05
            if csL.background and colH_len_L.inches > min_draw_in:
                bgL = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y, col_w_len, colH_len_L)
                bgL.fill.solid(); bgL.fill.fore_color.rgb = csL.background; bgL.line.fill.background()
            if csR.background and colH_len_R.inches > min_draw_in:
                bgR = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y, col_w_len, colH_len_R)
                bgR.fill.solid(); bgR.fill.fore_color.rgb = csR.background; bgR.line.fill.background()

            if csL.border_left_color and csL.border_left_width:
                bL = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y, csL.border_left_width, colH_len_L)
                bL.fill.solid(); bL.fill.fore_color.rgb = csL.border_left_color; bL.line.fill.background()
            if csR.border_left_color and csR.border_left_width:
                bR = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y, csR.border_left_width, colH_len_R)
                bR.fill.solid(); bR.fill.fore_color.rgb = csR.border_left_color; bR.line.fill.background()

            def render_into(x_left_len, y_top_len, avail_w_len, avail_w_in, elements_in_col):
                y_col = y_top_len
                base_pt = 18 * scale_factor
                for b in elements_in_col:
                    bt = b.get('type'); bs: StyleInfo = b.get('style') or StyleInfo()
                    if bt in ('h2','h3'):
                        fs = (bs.font_size.pt if bs.font_size else (28 if bt=='h2' else 22)) * scale_factor
                        h = _estimate_runs_height_in_inches(b.get('runs', []), fs, avail_w_in, (bs.line_height or 1.2))
                        box = slide.shapes.add_textbox(x_left_len, y_col, avail_w_len, h)
                        tf = box.text_frame; tf.word_wrap = True
                        p = tf.paragraphs[0]
                        for rd in b.get('runs', []) or [{'text': ''}]:
                            _emit_text_with_links_and_citations(
                                p,
                                text=rd.get('text', ''),
                                default_size_pt=fs,
                                bold=True,
                                italic=bool(rd.get('italic', False)),
                                color=(rd.get('color') or bs.color or DEFAULT_COLORS.get('text')),
                                sources_map=sources_map,
                                resolve_citations=resolve_citations
                            )
                        y_col += h + Inches(0.06 * scale_factor)

                    elif bt == 'paragraph':
                        fs = (bs.font_size.pt if bs.font_size else 18) * scale_factor
                        h = _estimate_runs_height_in_inches(b.get('runs', []), fs, avail_w_in, (bs.line_height or 1.3))
                        box = slide.shapes.add_textbox(x_left_len, y_col, avail_w_len, h)
                        tf = box.text_frame; tf.word_wrap = True
                        p = tf.paragraphs[0]
                        _emit_runs_with_sources(
                            p,
                            b.get('runs', []),
                            default_size_pt=fs,
                            fallback_color=(bs.color or DEFAULT_COLORS.get('text')),
                            sources_map=sources_map,
                            resolve_citations=resolve_citations
                        )
                        y_col += h + Inches(0.04 * scale_factor)

                    elif bt == 'list':
                        fs = base_pt
                        w_len = avail_w_len - Inches(0.25)
                        w_in  = max(0.1, avail_w_in - 0.25)
                        for it in b.get('items', []):
                            h = _estimate_runs_height_in_inches(it.get('runs', []), fs, w_in, 1.3)
                            box = slide.shapes.add_textbox(x_left_len + Inches(0.25), y_col, w_len, h)
                            tf = box.text_frame; tf.word_wrap = True
                            p = tf.paragraphs[0]; p.level = 0
                            _emit_runs_with_sources(
                                p,
                                it.get('runs', []),
                                default_size_pt=fs,
                                fallback_color=DEFAULT_COLORS.get('text'),
                                sources_map=sources_map,
                                resolve_citations=resolve_citations
                            )
                            y_col += h + Inches(0.04 * scale_factor)
                        y_col += Inches(0.02 * scale_factor)

                    elif bt == 'callout':
                        fs = (bs.font_size.pt if bs.font_size else 16) * scale_factor
                        pad_l = bs.padding_left or Inches(0.2)
                        pad_r = bs.padding_right or Inches(0.1)
                        pad_t = bs.padding_top or Inches(0.1)
                        pad_b = bs.padding_bottom or Inches(0.1)
                        inner_w_in  = max(0.1, avail_w_in - pad_l.inches - pad_r.inches)
                        inner_w_len = avail_w_len - pad_l - pad_r
                        content_h = _estimate_runs_height_in_inches(b.get('runs', []), fs, inner_w_in, (bs.line_height or 1.3))
                        box_h = Inches(max(content_h.inches + pad_t.inches + pad_b.inches, 0.6))

                        if bs.background:
                            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_left_len, y_col, avail_w_len, box_h)
                            bg.fill.solid(); bg.fill.fore_color.rgb = bs.background; bg.line.fill.background()
                        if bs.border_left_color and bs.border_left_width:
                            border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_left_len, y_col, bs.border_left_width, box_h)
                            border.fill.solid(); border.fill.fore_color.rgb = bs.border_left_color; border.line.fill.background()

                        tb = slide.shapes.add_textbox(x_left_len + pad_l, y_col + pad_t, inner_w_len, box_h - pad_t - pad_b)
                        tf = tb.text_frame; tf.word_wrap = True
                        p = tf.paragraphs[0]
                        _emit_runs_with_sources(
                            p,
                            b.get('runs', []),
                            default_size_pt=fs,
                            fallback_color=(bs.color or DEFAULT_COLORS.get('text')),
                            sources_map=sources_map,
                            resolve_citations=resolve_citations
                        )
                        y_col += box_h + Inches(0.08 * scale_factor)

                    elif bt == 'table':
                        header = b.get('header', [])
                        rows = b.get('rows', [])

                        if not header and rows:
                            header, rows = rows[0], rows[1:]
                        if not header:
                            continue
                        ncols = len(header)

                        # target sizes (initial); actual shape height may differ
                        row_h = Inches(0.30 * scale_factor)
                        header_h = Inches(0.35 * scale_factor)

                        th_rules = css_parser.styles.get('th', {})
                        stripe_rules = css_parser.styles.get('tr:nth-child(even)', {})
                        parse_color = CSSParser()._parse_color
                        header_bg = parse_color(th_rules.get('background-color')) if th_rules else None
                        header_fg = parse_color(th_rules.get('color')) if th_rules else None
                        stripe_bg = parse_color(stripe_rules.get('background-color')) if stripe_rules else None

                        table_shape = slide.shapes.add_table(
                            1 + len(rows), ncols,
                            x_left_len, y_col,
                            avail_w_len, header_h + row_h * max(1, len(rows))
                        )
                        tbl = table_shape.table

                        # header
                        for c, cell in enumerate(header):
                            tc = tbl.cell(0, c)
                            tf = tc.text_frame; tf.clear()
                            p = tf.paragraphs[0]
                            _emit_runs_with_sources(
                                p,
                                cell.get('runs', []),
                                default_size_pt=14 * scale_factor,
                                fallback_color=(header_fg or DEFAULT_COLORS.get('text')),
                                sources_map=sources_map,
                                resolve_citations=resolve_citations
                            )
                            for rr in p.runs: rr.font.bold = True
                            if header_bg:
                                tc.fill.solid(); tc.fill.fore_color.rgb = header_bg

                        # body
                        for r, row in enumerate(rows, start=1):
                            for c, cell in enumerate(row):
                                tc = tbl.cell(r, c)
                                tf = tc.text_frame; tf.clear()
                                p = tf.paragraphs[0]
                                _emit_runs_with_sources(
                                    p,
                                    cell.get('runs', []),
                                    default_size_pt=14 * scale_factor,
                                    fallback_color=DEFAULT_COLORS.get('text'),
                                    sources_map=sources_map,
                                    resolve_citations=resolve_citations
                                )
                            if stripe_bg and (r % 2 == 0):
                                for c in range(ncols):
                                    tc = tbl.cell(r, c); tc.fill.solid(); tc.fill.fore_color.rgb = stripe_bg

                        # advance by the **actual** drawn height
                        y_col += table_shape.height + Inches(0.06 * scale_factor)

                return y_col

            yL_start = y + padL_t
            yR_start = y + padR_t
            xL_inset = left_x + padL_l
            xR_inset = right_x + padR_l
            wL_inner = col_w_len - padL_l - padL_r
            wR_inner = col_w_len - padR_l - padR_r

            y_left_end  = render_into(xL_inset, yL_start, wL_inner, innerL_in, colL.get('elements', []))
            y_right_end = render_into(xR_inset, yR_start, wR_inner, innerR_in, colR.get('elements', []))
            y += Inches(max(colH_len_L.inches, colH_len_R.inches)) + Inches(0.04 * scale_factor)

        elif elem_type == 'table':
            header = elem.get('header', [])
            rows = elem.get('rows', [])

            if not header and rows:
                header, rows = rows[0], rows[1:]
            if header:
                ncols = len(header)

                tbl_w_len = SLIDE_WIDTH - MARGIN*2
                row_h = Inches(0.35 * scale_factor)
                header_h = Inches(0.4 * scale_factor)

                th_rules = css_parser.styles.get('th', {})
                stripe_rules = css_parser.styles.get('tr:nth-child(even)', {})
                parse_color = CSSParser()._parse_color
                header_bg = parse_color(th_rules.get('background-color')) if th_rules else None
                header_fg = parse_color(th_rules.get('color')) if th_rules else None
                stripe_bg = parse_color(stripe_rules.get('background-color')) if stripe_rules else None

                table_shape = slide.shapes.add_table(
                    1 + len(rows), ncols, MARGIN, y,
                    tbl_w_len, header_h + row_h * max(1, len(rows))
                )
                tbl = table_shape.table

                # header cells
                for c, cell in enumerate(header):
                    tc = tbl.cell(0, c)
                    tf = tc.text_frame; tf.clear()
                    p = tf.paragraphs[0]
                    _emit_runs_with_sources(
                        p,
                        cell.get('runs', []),
                        default_size_pt=14 * scale_factor,
                        fallback_color=(header_fg or DEFAULT_COLORS.get('text')),
                        sources_map=sources_map,
                        resolve_citations=resolve_citations
                    )
                    for rr in p.runs: rr.font.bold = True
                    if header_bg:
                        tc.fill.solid(); tc.fill.fore_color.rgb = header_bg

                # body rows
                for r, row in enumerate(rows, start=1):
                    for c, cell in enumerate(row):
                        tc = tbl.cell(r, c)
                        tf = tc.text_frame; tf.clear()
                        p = tf.paragraphs[0]
                        _emit_runs_with_sources(
                            p,
                            cell.get('runs', []),
                            default_size_pt=14 * scale_factor,
                            fallback_color=DEFAULT_COLORS.get('text'),
                            sources_map=sources_map,
                            resolve_citations=resolve_citations
                        )
                    if stripe_bg and (r % 2 == 0):
                        for c in range(ncols):
                            tc = tbl.cell(r, c); tc.fill.solid(); tc.fill.fore_color.rgb = stripe_bg

                # advance by the **actual** table height so the next H2 never overlaps
                y += table_shape.height + Inches(0.12 * scale_factor)

        elif elem_type == 'paragraph':
            fs = (style.font_size.pt if style.font_size else 18) * scale_factor
            h = _estimate_runs_height_in_inches(runs, fs, page_w_in, (style.line_height or 1.3))
            box = slide.shapes.add_textbox(MARGIN, y, SLIDE_WIDTH - 2*MARGIN, h)
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            _emit_runs_with_sources(
                p,
                runs,
                default_size_pt=fs,
                fallback_color=(style.color or DEFAULT_COLORS.get('text')),
                sources_map=sources_map,
                resolve_citations=resolve_citations
            )
            y += h + Inches(0.08 * scale_factor)


# === Sources / citations helpers ============================================

_LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
_CIT_TOKEN_RE = re.compile(r"\[\[S:([0-9,\s]+)\]\]")

def _domain_of(url: str) -> str:
    try:
        from urllib.parse import urlparse
        net = urlparse(url).netloc
        return net or url
    except Exception:
        return url

def _normalize_html_citations(html: str) -> str:
    """
    Convert <sup class='cite'>…</sup> patterns into [[S:...]] tokens so we can process consistently.
    Examples handled:
      <sup class="cite" data-sids="1,3">[S:1,3]</sup> -> [[S:1,3]]
      <sup class="cite">[S:2]</sup>                   -> [[S:2]]
    """
    html = re.sub(
        r'<sup\s+[^>]*class=["\'][^"\']*\bcite\b[^"\']*["\'][^>]*>\s*\[S:([^\]]+)\]\s*</sup>',
        r'[[S:\1]]',
        html,
        flags=re.I
    )
    # Fallback: strip any leftover <sup class="cite">…</sup> tags but keep their inner text
    html = re.sub(
        r'<sup\s+[^>]*class=["\'][^"\']*\bcite\b[^"\']*["\'][^>]*>([^<]*)</sup>',
        r'\1',
        html,
        flags=re.I
    )
    return html

def _emit_text_with_links_and_citations(
        paragraph,
        *,
        text: str,
        default_size_pt: float,
        bold: bool,
        italic: bool,
        color: Optional[RGBColor],
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool
):
    """
    Split plain text for [text](url) links and [[S:n,n2]] citation tokens,
    emit as multiple runs with appropriate hyperlink + accent color.

    IMPORTANT: For citations we render concise anchors "[n]" (never the title),
    to avoid inflating the line length compared to the original [[S:n]] token.
    """
    accent = DEFAULT_COLORS.get('primary', RGBColor(0, 102, 204))

    def _emit_run(t, col=color, link=None, b=bold, i=italic):
        if not t:
            return
        r = paragraph.add_run()
        r.text = t
        r.font.size = Pt(default_size_pt)
        r.font.bold = b
        r.font.italic = i
        if col:
            r.font.color.rgb = col
        if link:
            try:
                r.hyperlink.address = link
            except Exception:
                pass

    remaining = text
    while remaining:
        m_link = _LINK_RE.search(remaining)
        m_cit  = _CIT_TOKEN_RE.search(remaining) if resolve_citations else None

        # choose earliest match
        candidates = []
        if m_link: candidates.append(("link", m_link))
        if m_cit:  candidates.append(("cit",  m_cit))

        if not candidates:
            _emit_run(remaining)
            break

        kind, m = min(candidates, key=lambda x: x[1].start())

        # pre-text
        if m.start() > 0:
            _emit_run(remaining[:m.start()])

        if kind == "link":
            label, url = m.group(1), m.group(2)
            _emit_run(label, col=accent, link=url)
        else:
            # [[S:1, 2, 3]] → clickable "[1] · [2] · [3]"
            raw_ids = m.group(1)
            ids = []
            for part in raw_ids.split(','):
                part = part.strip()
                if part.isdigit():
                    ids.append(int(part))
            for idx, sid in enumerate(ids):
                rec = sources_map.get(sid, {})
                url = rec.get("url", "")
                _emit_run(f"[{sid}]", col=accent, link=(url or None))
                if idx != len(ids) - 1:
                    _emit_run(" · ")  # small separator between multiple cites

        # advance
        remaining = remaining[m.end():]

def _emit_runs_with_sources(
        paragraph,
        runs: List[Dict[str, Any]],
        *,
        default_size_pt: float,
        fallback_color: Optional[RGBColor],
        sources_map: Dict[int, Dict[str, str]],
        resolve_citations: bool
):
    """
    Emit a list of run dicts (as produced by our parser) while resolving links and citations.
    """
    for rd in runs or []:
        txt = rd.get('text') or ''
        if not txt:
            continue
        bold = rd.get('bold', False)
        italic = rd.get('italic', False)
        col = rd.get('color') or fallback_color or DEFAULT_COLORS.get('text')

        _emit_text_with_links_and_citations(
            paragraph,
            text=txt,
            default_size_pt=default_size_pt,
            bold=bold,
            italic=italic,
            color=col,
            sources_map=sources_map,
            resolve_citations=resolve_citations
        )

def _add_sources_slide(prs, sources_map: Dict[int, Dict[str, str]], order: List[int]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(0.6))
    tp = title_box.text_frame.paragraphs[0]
    r = tp.add_run(); r.text = "Sources"
    r.font.size = Pt(26); r.font.bold = True

    # Body list
    y = Inches(1.4)
    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(8.4), Inches(5.5))
    tf = box.text_frame
    p = tf.paragraphs[0]

    for i, sid in enumerate(order):
        if i > 0:
            p = tf.add_paragraph()
        src = sources_map.get(sid, {})
        title = src.get("title", f"Source {sid}")
        url = src.get("url", "")

        # [n] Title (domain)
        label = f"[{sid}] {title}"
        dr = p.add_run(); dr.text = label
        dr.font.size = Pt(16); dr.font.bold = False

        if url:
            dom = _domain_of(url)
            pr = p.add_run(); pr.text = f" ({dom})"
            pr.font.size = Pt(12)
            pr.font.color.rgb = DEFAULT_COLORS.get('primary', RGBColor(0, 102, 204))
            try:
                pr.hyperlink.address = url
            except Exception:
                pass

# ============================================================================
# MAIN
# ============================================================================

def render_pptx(
        path: str,
        content_md: str = "",
        content_html: str = "",
        *,
        title: Optional[str] = None,
        base_dir: Optional[str] = None,
        sources: Optional[str] = None,
        resolve_citations: bool = False,
        include_sources_slide: bool = False
) -> str:
    """Render PowerPoint from HTML."""

    outdir = resolve_output_dir()
    filename = Path(path).name
    if not filename.endswith('.pptx'):
        filename += '.pptx'
    outfile = outdir / filename
    outfile.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    if content_html and content_html.strip():
        # Parse CSS
        css_parser = CSSParser()
        style_match = re.search(r'<style[^>]*>(.*?)</style>', content_html, re.DOTALL | re.I)
        if style_match:
            css_parser.parse(style_match.group(1))

            # Debug: print parsed styles
            print("=== PARSED CSS ===")
            for selector, rules in css_parser.styles.items():
                print(f"{selector}: {rules}")

        # --- sources map
        sources_map: Dict[int, Dict[str, str]] = {}
        order: List[int] = []
        if sources:
            try:
                sources_map, order = md_utils._normalize_sources(sources)
            except Exception:
                # very defensive fallback
                try:
                    raw = json.loads(sources)
                    if isinstance(raw, dict):
                        sources_map = {int(k): v for k, v in raw.items() if str(k).isdigit()}
                        order = sorted(sources_map.keys())
                    elif isinstance(raw, list):
                        for rec in raw:
                            sid = int(rec.get("id"))
                            sources_map[sid] = {"title": rec.get("title",""), "url": rec.get("url","")}
                        order = sorted(sources_map.keys())
                except Exception:
                    pass

        # normalize <sup class='cite'> → [[S:...]]
        normalized_html = _normalize_html_citations(content_html)

        # Parse HTML
        html_parser = StyledHTMLParser(css_parser)
        html_parser.feed(normalized_html)

        slides_data = html_parser.slides
        for sidx, s in enumerate(slides_data, 1):
            for e in s.get('elements', []):
                if not e:
                    continue
                if e.get('type') == 'columns':
                    cols = e.get('columns', [])
                    l = len(cols[0].get('elements', [])) if len(cols) > 0 else 0
                    r = len(cols[1].get('elements', [])) if len(cols) > 1 else 0
                    print(f"Slide {sidx}: COLUMNS: {len(cols)} L:{l} R:{r}")

        print(f"=== PARSED {len(slides_data)} SLIDES ===")

        # Render ALL slides
        for slide_data in slides_data:
            print(f"Rendering slide: {slide_data.get('title')}")
            render_slide(prs, slide_data, css_parser, sources_map, resolve_citations)

        # Always add Sources slide when we have sources (per spec)
        if sources_map:
            _add_sources_slide(prs, sources_map, order)

    prs.save(str(outfile))
    return filename